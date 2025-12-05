"""
Utility helpers for common PowerPoint edits using python-pptx.
These wrap risky/verbose operations into safe, tested functions to reduce model hallucinations.
"""
from typing import List, Tuple, Optional, Any, Union

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR, MSO_SHAPE_TYPE
from pptx.oxml.xmlchemy import OxmlElement
from pptx.shapes.base import BaseShape
from pptx.slide import Slide


# -------- Basic getters --------
def load_presentation(ppt_path: str) -> Presentation:
    """Load presentation from path."""
    return Presentation(ppt_path)


def get_slide(prs: Presentation, index: int) -> Slide:
    """Safe slide access with 0-based index."""
    if index < 0 or index >= len(prs.slides):
        raise IndexError(f"Slide index {index} out of range; total slides: {len(prs.slides)}")
    return prs.slides[index]


# -------- Shape utilities --------
def delete_shapes_except(slide: Slide, shapes_to_keep: List[Optional[BaseShape]]) -> None:
    """Delete all shapes except those in shapes_to_keep."""
    keep_elements = {s.element for s in shapes_to_keep if s is not None}
    for shape in list(slide.shapes):
        if shape.element in keep_elements:
            continue
        sp = shape.element
        sp.getparent().remove(sp)


def remove_connectors_and_lines(slide: Slide) -> None:
    """Remove existing connector/line shapes."""
    connector_type = getattr(MSO_SHAPE_TYPE, "CONNECTOR", None)
    for shape in list(slide.shapes):
        if shape.shape_type == MSO_SHAPE_TYPE.LINE or (connector_type and shape.shape_type == connector_type):
            sp = shape.element
            sp.getparent().remove(sp)


def add_rounded_textbox(
    slide: Slide,
    text: str,
    left: Union[int, Emu],
    top: Union[int, Emu],
    width: Union[int, Emu],
    height: Union[int, Emu],
    fill_rgb: Tuple[int, int, int] = (232, 244, 248),
    text_rgb: Tuple[int, int, int] = (50, 50, 50),
    font_size: int = 20,
    align: Any = PP_ALIGN.CENTER,
    vertical_anchor: Any = MSO_ANCHOR.MIDDLE,
) -> BaseShape:
    """Add a rounded rectangle with text and return the shape."""
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height)

    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(*fill_rgb)

    line = shape.line
    line.color.rgb = RGBColor(180, 180, 180)
    line.width = Pt(1)

    if hasattr(shape, "adjustments") and len(shape.adjustments) > 0:
        shape.adjustments[0] = 0.15  # rounded corner

    # Shadow (color not exposed in python-pptx; leave default)
    shadow = shape.shadow
    shadow.inherit = False
    shadow.style = "OUTER"
    shadow.distance = Pt(3)
    shadow.angle = 45
    shadow.blur_radius = Pt(4)
    shadow.transparency = 0.5

    tf = shape.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = vertical_anchor
    tf.margin_left = tf.margin_right = Inches(0.1)
    tf.margin_top = tf.margin_bottom = Inches(0.1)

    p = tf.paragraphs[0]
    p.text = text
    p.font.name = "Microsoft JhengHei"
    p.font.size = Pt(font_size)
    p.font.color.rgb = RGBColor(*text_rgb)
    p.alignment = align

    return shape


def add_arrow_between(
    slide: Slide,
    shape_from: BaseShape,
    shape_to: BaseShape,
    color_rgb: Tuple[int, int, int] = (70, 70, 70),
    width_pt: float = 2.5,
    arrow_head: int = 2,
) -> BaseShape:
    """Add straight connector arrow between two shapes and return the connector."""
    start_x = Emu(shape_from.left + shape_from.width)
    start_y = Emu(shape_from.top + shape_from.height // 2)
    end_x = Emu(shape_to.left)
    end_y = Emu(shape_to.top + shape_to.height // 2)

    connector = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, start_x, start_y, end_x, end_y)
    connector.line.width = Pt(width_pt)
    connector.line.color.rgb = RGBColor(*color_rgb)
    # Apply a visible arrowhead via XML since python-pptx lacks arrow enums in older versions
    ln = connector._element.xpath(".//a:ln")
    if ln:
        tail_end = OxmlElement("a:tailEnd")
        tail_end.set("type", "triangle")
        ln[0].append(tail_end)
    return connector


# -------- Layout helpers --------
def distribute_horizontally(
    slide_width: Union[int, Emu],
    count: int,
    box_width: Union[int, Emu],
    gap: Union[int, Emu],
    margin: Union[int, Emu] = Inches(0.5),
) -> List[int]:
    """
    Compute left positions to distribute `count` boxes horizontally.
    Returns list of left positions (Emu).
    """
    total_width = count * box_width + (count - 1) * gap
    start_left = (slide_width - total_width) // 2
    return [start_left + i * (box_width + gap) for i in range(count)]
