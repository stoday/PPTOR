import os
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

class PPTInspector:
    def __init__(self, converter):
        self.converter = converter

    def inspect(self, ppt_path: str):
        """
        讀取 PPT 檔案，回傳：
        1. 文字結構摘要 (str)
        2. PDF 視覺檔案路徑 (str or None)
        """
        if not os.path.exists(ppt_path):
            return "File not found.", None

        # 1. 讀取文字結構
        text_summary = self._get_text_summary(ppt_path)

        # 2. 轉換為 PDF (視覺)
        pdf_path = self.converter.convert_to_pdf(ppt_path, output_dir="./temp_visuals")

        return text_summary, pdf_path

    def _get_text_summary(self, ppt_path):
        prs = Presentation(ppt_path)
        summary = []
        summary.append(f"Presentation: {os.path.basename(ppt_path)}")
        summary.append(f"Total Slides: {len(prs.slides)}")
        
        for i, slide in enumerate(prs.slides):
            slide_info = [f"\n--- Slide {i+1} ---"]
            
            # Layout
            slide_info.append(f"Layout: {slide.slide_layout.name}")
            
            # Title
            if slide.shapes.title and slide.shapes.title.text:
                slide_info.append(f"Title: '{slide.shapes.title.text}'")
            
            # Elements
            elements = []
            for shape in slide.shapes:
                # Skip title as it's already handled
                if shape == slide.shapes.title:
                    continue
                
                shape_desc = f"- Type: {shape.shape_type}"
                
                # Text content
                if shape.has_text_frame and shape.text.strip():
                    text = shape.text.replace('\n', ' ')
                    if len(text) > 50: text = text[:50] + "..."
                    shape_desc += f", Text: '{text}'"
                
                # Geometry
                try:
                    left_pt = int(float(shape.left.pt))
                    top_pt = int(float(shape.top.pt))
                    width_pt = int(float(shape.width.pt))
                    height_pt = int(float(shape.height.pt))
                    shape_desc += f", Pos: ({left_pt}, {top_pt}), Size: {width_pt}x{height_pt}"
                except Exception:
                    shape_desc += ", Pos/Size: (unreadable)"
                
                # Color (Simplified) - 安全地讀取顏色
                try:
                    if hasattr(shape, 'fill'):
                        fill = shape.fill
                        # 只有 solid 或 patterned 填充才有 fore_color
                        if fill.type in [1, 2]:  # 1=SOLID, 2=PATTERNED
                            if hasattr(fill, 'fore_color') and hasattr(fill.fore_color, 'rgb'):
                                shape_desc += f", Fill: #{fill.fore_color.rgb}"
                except Exception:
                    # 忽略無法讀取顏色的情況
                    pass

                elements.append(shape_desc)
            
            if elements:
                slide_info.append("Elements:")
                slide_info.extend(elements)
            else:
                slide_info.append("Elements: (None)")
                
            summary.append("\n".join(slide_info))
            
        return "\n".join(summary)
