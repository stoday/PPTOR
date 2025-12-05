import os
import re
import google.generativeai as genai
from pptx import Presentation
from pathlib import Path
from ppt_tool.converter import PPTConverter

# 設定 API Key
# 請確保環境變數 GOOGLE_API_KEY 已設定
if "GOOGLE_API_KEY" not in os.environ:
    print("[WARN] Warning: GOOGLE_API_KEY not found in environment variables.")

class PPTModifier:
    def __init__(self):
        try:
            genai.configure(api_key=os.environ.get("GOOGLE_API_KEY"))
            
            # 從環境變數讀取模型設定，若無則使用預設值
            self.text_model_name = os.environ.get("GEMINI_TEXT_MODEL", "gemini-2.5-flash")
            self.vision_model_name = os.environ.get("GEMINI_VISION_MODEL", "gemini-2.5-flash")
            
            # 初始化兩個模型
            self.text_model = genai.GenerativeModel(self.text_model_name)
            self.vision_model = genai.GenerativeModel(self.vision_model_name)
            
            print(f"[INFO] Text model: {self.text_model_name}")
            print(f"[INFO] Vision model: {self.vision_model_name}")
        except Exception as e:
            print(f"[WARN] Gemini Init Error: {e}")
            self.text_model = None
            self.vision_model = None

    def generate_and_execute(self, user_instruction, text_summary, pdf_path, ppt_path, debug=False):
        """
        1. 呼叫 Gemini 生成程式碼
        2. 執行程式碼修改 PPT
        """
        if not self.text_model or not self.vision_model:
            return False, "Gemini model not initialized."

        # 確保生成的程式碼可以直接取用目標檔案路徑
        exec_globals = {"ppt_path": ppt_path, "__builtins__": __builtins__}
        code = ""

        # 如果 PowerPoint 正在開啟會留下 lock 檔 (~$filename)，先提醒用戶關閉
        if self._has_lock_file(ppt_path):
            return False, "PowerPoint 似乎正在開啟檔案（偵測到 ~$ 開頭的 lock 檔）。請先關閉簡報再重試。"

        # 根據是否有 PDF 選擇模型
        if pdf_path and os.path.exists(pdf_path):
            model = self.vision_model
            model_name = self.vision_model_name
            print(f"[INFO] Using vision model ({model_name}) for PDF analysis...")
        else:
            model = self.text_model
            model_name = self.text_model_name
            print(f"[INFO] Using text model ({model_name})...")
        
        print("[INFO] Building prompt for Gemini...")
        # 準備 Prompt
        prompt_parts = [
            "You are an expert Python developer specializing in `python-pptx`.",
            "Your task is to write a Python script to modify an existing PowerPoint file based on the user's instruction and the current state of the presentation.",
            f"The target PowerPoint file path is: `{ppt_path}`",
            "The user does NOT provide code; you must choose the correct helper APIs yourself (see below).",
            "You MUST use `Presentation(ppt_path)` to load the file, make changes, and then `prs.save(ppt_path)` to save it.",
            "\n--- Current Presentation Structure ---",
            text_summary,
            "\n--- User Instruction ---",
            user_instruction,
            "\n--- CRITICAL Requirements ---",
            "1. Output ONLY valid Python code. No markdown fences, no explanations.",
            "2. Always import required modules at the top:",
            "   from pptx import Presentation",
            "   from pptx.util import Inches, Pt",
            "   from pptx.dml.color import RGBColor",
            "   from pptx.enum.text import PP_ALIGN",
            "   from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR, MSO_SHAPE_TYPE",
            "   from pptx.enum.text import MSO_ANCHOR  # For vertical centering, use MSO_ANCHOR",
            "   from ppt_tool.ppt_api import (",
            "       load_presentation, get_slide, delete_shapes_except, remove_connectors_and_lines,",
            "       add_rounded_textbox, add_arrow_between, distribute_horizontally",
            "   )",
            "   # You MUST call these helper APIs instead of raw slide.shapes.add_* unless a needed helper is missing.",
            "",
            "   - Do NOT swallow errors when saving: avoid wrapping `prs.save` in try/except; if you must catch, re-raise the exception so the caller sees the failure.",
            "",
            "3. RGBColor MUST have THREE arguments: RGBColor(r, g, b) where each is 0-255",
            "   CORRECT: RGBColor(255, 0, 0)  # Red",
            "   WRONG: RGBColor(0xFF0000)",
            "",
            "4. To DELETE old elements from a slide:",
            "   # Delete all shapes except title",
            "   shapes_to_delete = [shape for shape in slide.shapes if shape != slide.shapes.title]",
            "   for shape in shapes_to_delete:",
            "       sp = shape.element",
            "       sp.getparent().remove(sp)",
            "",
            "5. To calculate proper textbox size for Chinese text:",
            "   - Width: at least 15 characters * Pt(font_size) * 0.6",
            "   - Height: at least Pt(font_size) * 1.5 * number_of_lines",
            "   - Add padding: width += Inches(0.2), height += Inches(0.1)",
            "",
            "6. For horizontal layout of N boxes:",
            "   slide_width = Inches(10)  # Standard 16:9 slide",
            "   box_width = (slide_width - Inches(1)) / N  # Leave margins",
            "   for i in range(N):",
            "       left = Inches(0.5) + i * box_width",
            "",
            "7. Text formatting (CRITICAL for textboxes):",
            "   - ALWAYS enable word wrap: text_frame.word_wrap = True",
            "   - Chinese font: paragraph.font.name = 'Microsoft JhengHei'",
            "   - Text color: paragraph.font.color.rgb = RGBColor(r, g, b)",
            "   - Alignment: paragraph.alignment = PP_ALIGN.CENTER (or LEFT, RIGHT)",
            "   - WRONG: PP_PARAGRAPH_ALIGNMENT.MIDDLE / PP_ALIGN.MIDDLE (do not exist!)",
            "   - CORRECT: PP_ALIGN.CENTER",
            "   - Vertical centering: text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE (requires `from pptx.enum.text import MSO_ANCHOR`)",
            "   - Text detection: use `shape.has_text_frame` and check `shape.text_frame.text.strip()`; there is NO `text_frame.has_text` attribute.",
            "   - Without word_wrap=True, text will overflow the box!",
            "",
            "8. Design aesthetics (make it beautiful!):",
            "   - Use harmonious color palettes (e.g., blue theme: #4A90E2, #50E3C2, #E8F4F8)",
            "   - Add rounded corners to shapes: shape.adjustments[0] = 0.1",
            "   - Add subtle shadows for depth:",
            "     shape.shadow.inherit = False",
            "     shape.shadow.style = 'OUTER'",  # color/fore_color not exposed; do not set",
            "     shape.shadow.distance = Pt(3)",
            "     shape.shadow.angle = 45",
            "     shape.shadow.blur_radius = Pt(4)",
            "     shape.shadow.transparency = 0.5",
            "   - Use consistent spacing and alignment",
            "   - Prefer soft colors over harsh primary colors",
            "   - Add shape fills with light backgrounds: RGBColor(232, 244, 248) for blue theme",
            "   - Rounded corners safety: only set `shape.adjustments[0]` if `hasattr(shape, 'adjustments') and len(shape.adjustments)>0`; otherwise use `slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, ...)` and write text in its text_frame.",
            "   - Shadow safety: ShadowFormat in python-pptx has no `color`/`fore_color`; do NOT set those properties. Just set inherit/style/distance/angle/blur_radius/transparency.",
            "",
            "9. Drawing arrows/connectors:",
            "   from pptx.enum.shapes import MSO_CONNECTOR",
            "   from pptx.util import Inches",
            "   # Draw arrow from (x1,y1) to (x2,y2)",
            "   connector = slide.shapes.add_connector(",
            "       MSO_CONNECTOR.STRAIGHT,",
            "       Inches(x1), Inches(y1),",
            "       Inches(x2), Inches(y2)",
            "   )",
            "   connector.line.width = Pt(2)",
            "   # Add arrowhead at end",
            "   connector.line.end_arrow_type = 2  # 2 = arrow",
            "   # WRONG: MSO_ARROWHEAD (does not exist!)",
            "   # CORRECT: Use connector.line.end_arrow_type = 2",
            "",
            "10. Safe indexing (CRITICAL - avoid index errors!):",
            "   # ALWAYS check bounds before accessing slides",
            "   if len(prs.slides) > 5:  # Check before accessing slides[5]",
            "       slide = prs.slides[5]",
            "   # For 'page 6', use index 5 (0-based)",
            "   # WRONG: slide = prs.slides[6]  # This is the 7th slide!",
            "   # CORRECT: slide = prs.slides[5]  # This is the 6th slide",
            "",
            "11. Error handling:",
            "   # Wrap risky operations in try-except",
            "   try:",
            "       # Your code here",
            "   except IndexError as e:",
            "       print(f'Index error: {e}')",
            "   except Exception as e:",
            "       print(f'Error: {e}')",
            "",
            "12. Shape type checks:",
            "   - For lines/connectors detection: use `shape.shape_type == MSO_SHAPE_TYPE.LINE`.",
            "   - Do NOT use `MSO_AUTO_SHAPE_TYPE.LINE` (does not exist).",
            "   - Auto shapes like rectangles: `shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE` or add specific shapes via MSO_AUTO_SHAPE_TYPE.",
            "",
            "13. Preferred helper API (use these instead of raw pptx calls when possible):",
            "   - load_presentation(ppt_path) -> prs",
            "   - get_slide(prs, index) -> slide",
            "   - delete_shapes_except(slide, [shapes_to_keep])",
            "   - remove_connectors_and_lines(slide)",
            "   - distribute_horizontally(slide_width, count, box_width, gap, margin=Inches(0.5)) -> list of left positions",
            "   - add_rounded_textbox(slide, text, left, top, width, height, fill_rgb=(232,244,248), text_rgb=(50,50,50), font_size=20)",
            "   - add_arrow_between(slide, shape_from, shape_to, color_rgb=(70,70,70), width_pt=2.5, arrow_head=2)",
            "   - ALWAYS prefer these helpers; avoid direct slide.shapes.add_* unless helper unavailable.",
            "",
            "\n--- Code Example: Replace bullet list with horizontal textboxes ---",
            "```python",
            "from pptx import Presentation",
            "from ppt_tool.ppt_api import (",
            "    load_presentation, get_slide, delete_shapes_except, remove_connectors_and_lines,",
            "    add_rounded_textbox, add_arrow_between, distribute_horizontally",
            ")",
            "from pptx.util import Inches, Pt",
            "from pptx.dml.color import RGBColor",
            "from pptx.enum.text import PP_ALIGN, MSO_ANCHOR",
            "",
            "prs = load_presentation('path.pptx')",
            "slide = get_slide(prs, 5)  # 6th slide (0-indexed)",
            "delete_shapes_except(slide, [slide.shapes.title])",
            "",
            "items = ['項目一', '項目二', '項目三']",
            "lefts = distribute_horizontally(prs.slide_width, len(items), Inches(3), Inches(0.2))",
            "boxes = []",
            "for i, text in enumerate(items):",
            "    box = add_rounded_textbox(slide, text, lefts[i], Inches(2.5), Inches(3), Inches(2),",
            "        fill_rgb=(232,244,248), text_rgb=(50,50,50), font_size=18, align=PP_ALIGN.CENTER, vertical_anchor=MSO_ANCHOR.MIDDLE)",
            "    boxes.append(box)",
            "remove_connectors_and_lines(slide)",
            "for i in range(len(boxes)-1):",
            "    add_arrow_between(slide, boxes[i], boxes[i+1], color_rgb=(70,70,70), width_pt=2.5)",
            "",
            "prs.save('path.pptx')",
            "```"
        ]

        # 如果有 PDF 視覺檔，上傳並加入 Prompt
        if pdf_path and os.path.exists(pdf_path):
            print(f"[INFO] Uploading visual reference: {pdf_path}")
            try:
                pdf_file = genai.upload_file(pdf_path)
                prompt_parts.append("I have attached a PDF rendering of the current slides for your visual reference. Use this to understand layout, alignment, and colors.")
                prompt_parts.append(pdf_file)
            except Exception as e:
                print(f"[WARN] Failed to upload PDF: {e}")

        # 呼叫 Gemini
        try:
            print("[INFO] Calling Gemini model...")
            response = model.generate_content(prompt_parts)
            code = self._extract_code(response.text)
            
            print("[INFO] Executing generated code...")
            if debug:
                print("\n[DEBUG] Generated code:\n" + "="*60)
                print(code)
                print("="*60 + "\n")

            # Enforce usage of helper API; reject obvious non-compliant code
            helper_markers = [
                "ppt_tool.ppt_api",
                "add_rounded_textbox",
                "add_arrow_between",
                "delete_shapes_except",
                "remove_connectors_and_lines",
                "distribute_horizontally",
            ]
            if not any(marker in code for marker in helper_markers):
                return False, "Generated code did not use required helper APIs; please retry."
            
            # 檢查檔案是否被鎖定（例如在 PowerPoint 中開啟）
            if self._is_file_locked(ppt_path):
                print("[WARN] PowerPoint file is currently open in another application.")
                print("[WARN] Please close PowerPoint and press Enter to continue...")
                input()
            
            # print(code) # Debug use
            
            # 執行程式碼
            # 為了安全，限制 globals，但允許 pptx 相關庫
            exec(code, exec_globals)
            
            return True, "Modification applied successfully."
            
        except Exception as e:
            # 輸出生成的程式碼供除錯
            print("\n[DEBUG] Generated code that caused error:")
            print("=" * 60)
            print(code)
            print("=" * 60)
            print("\n[DEBUG] Full traceback:")
            import traceback
            traceback.print_exc()
            return False, f"Error: {e}"
        
    def validate_with_vision(self, user_instruction: str, ppt_path: str):
        """
        使用視覺模型檢查修改後的簡報是否符合指令。
        回傳 (ok: bool, message: str)
        """
        if not self.vision_model:
            return False, "Vision model not initialized."
        
        print("[INFO] Validating result with vision model...")
        converter = PPTConverter()
        pdf_path = converter.convert_to_pdf(ppt_path, output_dir="./temp_visuals")
        if not pdf_path or not os.path.exists(pdf_path):
            return False, "Validation skipped: PDF conversion failed."
        
        try:
            print(f"[INFO] Uploading updated PDF for validation: {pdf_path}")
            pdf_file = genai.upload_file(pdf_path)
            prompt = [
                "You are a QA checker. Review the presentation PDF and compare against the user's instruction.",
                "List mismatches or missing elements; be concise.",
                "--- User instruction ---",
                user_instruction,
                "--- PDF attached below ---",
                pdf_file,
            ]
            resp = self.vision_model.generate_content(prompt)
            feedback = resp.text.strip() if resp and resp.text else "No feedback."
            print("[INFO] Validation result:")
            print(feedback)
            ok = "no issues" in feedback.lower() or "looks good" in feedback.lower()
            return ok, feedback
        except Exception as e:
            print(f"[WARN] Validation failed: {e}")
            return False, f"Validation error: {e}"

    def _is_file_locked(self, filepath):
        """檢查檔案是否被其他程式鎖定"""
        import os
        if not os.path.exists(filepath):
            return False
        try:
            # 嘗試以獨占模式開啟檔案
            with open(filepath, 'r+b') as f:
                pass
            return False
        except (IOError, PermissionError):
            return True

    def _has_lock_file(self, filepath):
        """偵測 PowerPoint 產生的 lock 檔 (~$xxx.pptx)"""
        path = Path(filepath)
        lock_name = "~$" + path.name
        lock_path = path.parent / lock_name
        return lock_path.exists()

    def _extract_code(self, text):
        """從回應中提取 Python 程式碼區塊"""
        match = re.search(r'```python\n(.*?)\n```', text, re.DOTALL)
        if match:
            return match.group(1)
        # 如果沒有 markdown 標記，假設全是程式碼
        return text
