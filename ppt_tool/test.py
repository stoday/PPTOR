from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE_TYPE
from ppt_tool.ppt_api import (
    load_presentation, get_slide, add_rounded_textbox
)

# Target PowerPoint file path
# NOTE: Update this path to your local presentation file
ppt_path = r'C:\Users\today\Dropbox\MainStorage\P2025_PPTOR\presentation.pptx'

# Load the presentation
prs = load_presentation(ppt_path)

# Get Slide 4 (index 3)
slide_index = 3
if slide_index >= len(prs.slides):
    raise IndexError(f"Slide at index {slide_index} does not exist.")
slide = get_slide(prs, slide_index)

# Define the detailed descriptions for each step
detailed_texts = [
    "乾性材料如麵粉、糖、泡打粉等，需要先用篩網過篩，確保沒有結塊，使蛋糕口感更細膩。混合時確保所有粉末均勻分佈，避免在烘烤時出現生粉團。",
    "濕性材料包括雞蛋、牛奶、植物油或融化的奶油、香草精等。先將雞蛋打散，再陸續加入牛奶、油和香草精，攪拌至乳化狀態。這些液體要充分混合，才能更好地與乾性材料結合。",
    "將濕性材料分三次左右加入乾性材料中，每次加入後用刮刀以「切拌」或「翻拌」的方式輕柔混合，直到沒有明顯的乾粉即可。避免過度攪拌，以免麵粉產生筋性，影響蛋糕的鬆軟度。輕柔攪拌是蛋糕成功的關鍵。"
]

# Find existing flowchart shapes to determine positioning for new description boxes
flowchart_boxes = []
for shape in slide.shapes:
    # Identify auto shapes that form the existing flowchart based on their text content
    if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE and shape.text_frame.has_text:
        text = shape.text_frame.text
        if '乾性材料' in text or '濕性材料' in text or '濕性材料分次加入' in text:
            flowchart_boxes.append(shape)

# Sort the identified boxes by their left position to match the order of detailed_texts
flowchart_boxes.sort(key=lambda s: s.left)

# Define common properties for the new detailed description boxes
new_box_height = Inches(3.0)  # Height adjusted to accommodate detailed text
# Place new boxes below the existing flowchart boxes.
# Existing boxes' top is Inches(180/72) = Inches(2.5), height is Inches(180/72) = Inches(2.5).
# So their bottom is at Inches(5.0). A gap of 0.5 inches makes new boxes start at 5.5 inches from top.
new_box_top = Inches(5.5)
fill_color_rgb = (232, 244, 248)  # Light blue, consistent with existing boxes
text_color_rgb = (50, 50, 50)     # Dark gray for readability
font_size_pt = 14                 # Slightly smaller font for detailed reference text

# Add new detailed description boxes below the corresponding flowchart boxes
for i, box_text in enumerate(detailed_texts):
    if i < len(flowchart_boxes):
        existing_box = flowchart_boxes[i]
        left = existing_box.left
        width = existing_box.width  # Use the same width as the existing flowchart boxes

        # Add the rounded textbox using the helper function
        new_description_box = add_rounded_textbox(
            slide, box_text, left, new_box_top, width, new_box_height,
            fill_rgb=fill_color_rgb, text_rgb=text_color_rgb, font_size=font_size_pt
        )

        # Apply specific text formatting as per critical requirements
        text_frame = new_description_box.text_frame
        text_frame.word_wrap = True  # Enable word wrap for detailed text
        text_frame.vertical_anchor = MSO_ANCHOR.TOP  # Align text to the top

        # Assuming the add_rounded_textbox helper places all text in the first paragraph
        p = text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT  # Left align for detailed text
        p.font.name = 'Microsoft JhengHei'  # Set Chinese font
        p.font.size = Pt(font_size_pt)  # Set font size
        p.font.color.rgb = RGBColor(*text_color_rgb)  # Set text color

        # Add subtle shadow for depth, as per critical requirements
        if hasattr(new_description_box, 'shadow') and new_description_box.shadow is not None:
            shadow = new_description_box.shadow
            shadow.inherit = False
            # Following the critical requirement for 'style' attribute
            shadow.style = 'OUTER'
            shadow.distance = Pt(3)
            shadow.angle = 45
            shadow.blur_radius = Pt(4)
            shadow.transparency = 0.5

# Save the modified presentation
prs.save(ppt_path)
