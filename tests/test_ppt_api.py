import pytest
from pathlib import Path

from ppt_tool.ppt_api import (
    load_presentation,
    get_slide,
    delete_shapes_except,
    remove_connectors_and_lines,
    add_rounded_textbox,
    add_arrow_between,
    distribute_horizontally,
)
from pptx import Presentation
from pptx.util import Inches
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.dml.color import RGBColor


@pytest.fixture(scope="session")
def artifacts_dir():
    path = Path("tests/artifacts")
    path.mkdir(parents=True, exist_ok=True)
    return path


def _new_single_slide_ppt(tmpdir: Path, title: str) -> Path:
    """Create a one-slide presentation with the given title."""
    ppt_path = tmpdir / f"{title}.pptx"
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])  # Title slide layout
    slide.shapes.title.text = title
    prs.save(ppt_path)
    return ppt_path


def _save_artifact(prs: Presentation, artifacts_dir: Path, name: str):
    out = artifacts_dir / name
    prs.save(out)
    return out


def test_load_and_get_slide(tmp_path: Path, artifacts_dir: Path):
    ppt_path = _new_single_slide_ppt(tmp_path, "load_and_get_slide")
    prs = load_presentation(str(ppt_path))
    slide = get_slide(prs, 0)
    assert slide.shapes.title.text == "load_and_get_slide"
    _save_artifact(prs, artifacts_dir, "load_and_get_slide.pptx")


def test_delete_shapes_except(tmp_path: Path, artifacts_dir: Path):
    ppt_path = _new_single_slide_ppt(tmp_path, "delete_shapes_except")
    prs = load_presentation(str(ppt_path))
    slide = get_slide(prs, 0)
    tb = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(2), Inches(1))
    tb.text = "temp"
    delete_shapes_except(slide, [slide.shapes.title])
    remaining_texts = [s.text for s in slide.shapes if hasattr(s, "text")]
    assert remaining_texts == ["delete_shapes_except"]
    _save_artifact(prs, artifacts_dir, "delete_shapes_except.pptx")


def test_remove_connectors_and_lines(tmp_path: Path, artifacts_dir: Path):
    ppt_path = _new_single_slide_ppt(tmp_path, "remove_connectors_and_lines")
    prs = load_presentation(str(ppt_path))
    slide = get_slide(prs, 0)
    s1 = add_rounded_textbox(slide, "A", Inches(1), Inches(3), Inches(2), Inches(1))
    s2 = add_rounded_textbox(slide, "B", Inches(4), Inches(3), Inches(2), Inches(1))
    add_arrow_between(slide, s1, s2)
    remove_connectors_and_lines(slide)
    connector_type = getattr(MSO_SHAPE_TYPE, "CONNECTOR", None)
    connectors = [
        s for s in slide.shapes
        if s.shape_type == MSO_SHAPE_TYPE.LINE or (connector_type and s.shape_type == connector_type)
    ]
    assert len(connectors) == 0
    _save_artifact(prs, artifacts_dir, "remove_connectors_and_lines.pptx")


def test_add_rounded_textbox(tmp_path: Path, artifacts_dir: Path):
    ppt_path = _new_single_slide_ppt(tmp_path, "add_rounded_textbox")
    prs = load_presentation(str(ppt_path))
    slide = get_slide(prs, 0)
    shape = add_rounded_textbox(slide, "Hello", Inches(1), Inches(3), Inches(3), Inches(1.5))
    assert shape.has_text_frame
    assert shape.text_frame.text == "Hello"
    assert shape.text_frame.word_wrap
    _save_artifact(prs, artifacts_dir, "add_rounded_textbox.pptx")


def test_add_arrow_between(tmp_path: Path, artifacts_dir: Path):
    ppt_path = _new_single_slide_ppt(tmp_path, "add_arrow_between")
    prs = load_presentation(str(ppt_path))
    slide = get_slide(prs, 0)
    lefts = distribute_horizontally(prs.slide_width, 2, Inches(2), Inches(1))
    s1 = add_rounded_textbox(slide, "Left", lefts[0], Inches(3), Inches(2), Inches(1.5))
    s2 = add_rounded_textbox(slide, "Right", lefts[1], Inches(3), Inches(2), Inches(1.5))
    before = len(slide.shapes)
    connector = add_arrow_between(slide, s1, s2, color_rgb=(10, 20, 30), width_pt=3, arrow_head=2)
    after = len(slide.shapes)
    assert after == before + 1
    assert connector.line.color.rgb == RGBColor(10, 20, 30)
    # Validate arrowhead exists in XML
    ln = connector._element.xpath(".//a:ln")
    tail = ln[0].find("{http://schemas.openxmlformats.org/drawingml/2006/main}tailEnd") if ln else None
    assert tail is not None
    assert tail.get("type") == "triangle"
    _save_artifact(prs, artifacts_dir, "add_arrow_between.pptx")


def test_distribute_horizontally(tmp_path: Path, artifacts_dir: Path):
    ppt_path = _new_single_slide_ppt(tmp_path, "distribute_horizontally")
    prs = load_presentation(str(ppt_path))
    positions = distribute_horizontally(prs.slide_width, 3, Inches(2), Inches(0.5))
    assert len(positions) == 3
    assert positions[0] < positions[1] < positions[2]
    span = positions[2] + Inches(2) - positions[0]
    assert span <= prs.slide_width
    _save_artifact(prs, artifacts_dir, "distribute_horizontally.pptx")
